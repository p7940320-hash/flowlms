from pymongo import MongoClient
import os
import uuid
from dotenv import load_dotenv
from datetime import datetime, timezone
from pptx import Presentation
from pptx.util import Inches, Pt
import io

load_dotenv()

MONGO_URL = os.getenv("MONGO_URL")
DB_NAME = os.getenv("DB_NAME")

client = MongoClient(MONGO_URL)
db = client[DB_NAME]

# PowerPoint files and their course details
ppt_files = [
    {
        "filename": "PRESENTATION ON PUMP SPARES.pptx",
        "title": "Pump Spares Presentation",
        "code": "ENG-PUMPS-001",
        "description": "Comprehensive presentation covering pump spare parts, maintenance components, and replacement procedures. Learn about essential spare parts inventory and proper handling.",
        "category": "Engineering"
    },
    {
        "filename": "Valve Presentation.pptx",
        "title": "Valve Presentation",
        "code": "ENG-VALVE-001",
        "description": "Complete guide to valve types, applications, and maintenance procedures. Learn about different valve configurations and selection criteria.",
        "category": "Engineering"
    },
    {
        "filename": "PUMP CURVES, SELECTION & MATERIAL CHOOSING.pptx",
        "title": "Pump Curves, Selection & Material Choosing",
        "code": "ENG-CURVES-001",
        "description": "Learn about pump performance curves, proper pump selection criteria, and material selection for different applications and operating conditions.",
        "category": "Engineering"
    },
    {
        "filename": "PUMP TYPES, PUMP PARTS & WORKING PRINCIPLES PRESENTATION EDITED.pptx",
        "title": "Pump Types, Parts & Working Principles",
        "code": "ENG-TYPES-001",
        "description": "Detailed presentation on various pump types, their components, and working principles. Understand the fundamentals of pump operation and design.",
        "category": "Engineering"
    }
]

uploads_dir = "c:\\Users\\User\\Desktop\\flowlms\\flowlms\\backend\\uploads"

for ppt_info in ppt_files:
    ppt_path = os.path.join(uploads_dir, ppt_info["filename"])
    
    if not os.path.exists(ppt_path):
        print(f"File not found: {ppt_info['filename']}")
        continue
    
    print(f"\nProcessing: {ppt_info['title']}")
    
    try:
        # Load PowerPoint
        prs = Presentation(ppt_path)
        total_slides = len(prs.slides)
        print(f"  Found {total_slides} slides")
        
        # Create course
        course_id = str(uuid.uuid4())
        course = {
            "id": course_id,
            "title": ppt_info["title"],
            "code": ppt_info["code"],
            "description": ppt_info["description"],
            "category": ppt_info["category"],
            "type": "optional",
            "duration": f"{max(1, total_slides // 10)} hour{'s' if total_slides > 10 else ''}",
            "thumbnail": "",
            "is_published": True,
            "enrolled_users": [],
            "createdAt": datetime.now(timezone.utc),
            "updatedAt": datetime.now(timezone.utc)
        }
        db.courses.insert_one(course)
        print(f"  Created course: {course['title']}")
        
        # Create module
        module_id = str(uuid.uuid4())
        module = {
            "id": module_id,
            "course_id": course_id,
            "title": "Course Content",
            "description": "",
            "order": 1
        }
        db.modules.insert_one(module)
        
        # Create lessons from slides
        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract text from slide
            slide_text = []
            slide_title = f"Slide {slide_num}"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Use first text as title if it's short
                    if slide_num == 1 or (len(text) < 100 and not slide_text):
                        slide_title = text[:100]
                    slide_text.append(text)
            
            # Create HTML content for the slide
            content_html = '<div class="lesson-content">'
            content_html += f'<div class="slide-page">'
            content_html += f'<div class="slide-number">Slide {slide_num} of {total_slides}</div>'
            
            if slide_text:
                for i, text in enumerate(slide_text):
                    if i == 0:
                        content_html += f'<h2>{text}</h2>'
                    elif len(text) > 200:
                        content_html += f'<p>{text}</p>'
                    else:
                        content_html += f'<h3>{text}</h3>'
            else:
                content_html += f'<p><em>Visual content - Slide {slide_num}</em></p>'
            
            content_html += '</div></div>'
            
            # Create lesson
            lesson_id = str(uuid.uuid4())
            lesson = {
                "id": lesson_id,
                "module_id": module_id,
                "title": slide_title if slide_title != f"Slide {slide_num}" else f"Slide {slide_num}",
                "content_type": "text",
                "content": content_html,
                "duration_minutes": 2,
                "order": slide_num
            }
            db.lessons.insert_one(lesson)
        
        print(f"  Created {total_slides} lessons")
        
    except Exception as e:
        print(f"  Error processing {ppt_info['filename']}: {str(e)}")
        continue

print("\n\nAll PowerPoint courses created successfully!")

# Show summary
courses = list(db.courses.find({}, {"_id": 0, "title": 1, "code": 1, "category": 1}))
print(f"\nTotal courses: {len(courses)}")
for course in courses:
    print(f"  - {course.get('title')} ({course.get('code')})")

client.close()
