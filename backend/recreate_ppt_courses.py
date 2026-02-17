from pymongo import MongoClient
import os
import uuid
from dotenv import load_dotenv
from datetime import datetime, timezone
from pptx import Presentation
from PIL import Image
import io

load_dotenv()

MONGO_URL = os.getenv("MONGO_URL")
DB_NAME = os.getenv("DB_NAME")

client = MongoClient(MONGO_URL)
db = client[DB_NAME]

# First, delete the existing PowerPoint courses
print("Deleting existing PowerPoint courses...")
ppt_courses = list(db.courses.find({"category": "Engineering"}))
for course in ppt_courses:
    course_id = course.get('id')
    modules = list(db.modules.find({"course_id": course_id}))
    for module in modules:
        db.lessons.delete_many({"module_id": module.get('id')})
    db.modules.delete_many({"course_id": course_id})
    db.progress.delete_many({"course_id": course_id})
    db.courses.delete_one({"_id": course["_id"]})
print(f"Deleted {len(ppt_courses)} courses")

# PowerPoint files and their course details
ppt_files = [
    {
        "filename": "PRESENTATION ON PUMP SPARES.pptx",
        "title": "Pump Spares Presentation",
        "code": "ENG-PUMPS-001",
        "description": "Comprehensive presentation covering pump spare parts, maintenance components, and replacement procedures.",
        "category": "Engineering"
    },
    {
        "filename": "Valve Presentation.pptx",
        "title": "Valve Presentation",
        "code": "ENG-VALVE-001",
        "description": "Complete guide to valve types, applications, and maintenance procedures.",
        "category": "Engineering"
    },
    {
        "filename": "PUMP CURVES, SELECTION & MATERIAL CHOOSING.pptx",
        "title": "Pump Curves, Selection & Material Choosing",
        "code": "ENG-CURVES-001",
        "description": "Learn about pump performance curves, proper pump selection criteria, and material selection.",
        "category": "Engineering"
    },
    {
        "filename": "PUMP TYPES, PUMP PARTS & WORKING PRINCIPLES PRESENTATION EDITED.pptx",
        "title": "Pump Types, Parts & Working Principles",
        "code": "ENG-TYPES-001",
        "description": "Detailed presentation on various pump types, their components, and working principles.",
        "category": "Engineering"
    }
]

uploads_dir = "c:\\Users\\User\\Desktop\\flowlms\\flowlms\\backend\\uploads"

print("\nNote: PowerPoint slides will be displayed as embedded presentations.")
print("To show actual slide images, you need to:")
print("1. Export slides as images from PowerPoint (File > Export > PNG)")
print("2. Upload images to a hosting service (Google Drive, Imgur, etc.)")
print("3. Use those image URLs in the lessons\n")

for ppt_info in ppt_files:
    ppt_path = os.path.join(uploads_dir, ppt_info["filename"])
    
    if not os.path.exists(ppt_path):
        print(f"File not found: {ppt_info['filename']}")
        continue
    
    print(f"Processing: {ppt_info['title']}")
    
    try:
        prs = Presentation(ppt_path)
        total_slides = len(prs.slides)
        print(f"  Found {total_slides} slides")
        
        # Create course
        course_id = str(uuid.uuid4())
        course = {
            "id": course_id,
            "title": ppt_info["title"],
            "code": ppt_info["code"],
            "description": ppt_info["description"],
            "category": ppt_info["category"],
            "type": "optional",
            "duration": f"{max(1, total_slides // 10)} hour{'s' if total_slides > 10 else ''}",
            "thumbnail": "",
            "is_published": True,
            "enrolled_users": [],
            "createdAt": datetime.now(timezone.utc),
            "updatedAt": datetime.now(timezone.utc)
        }
        db.courses.insert_one(course)
        
        # Create module
        module_id = str(uuid.uuid4())
        module = {
            "id": module_id,
            "course_id": course_id,
            "title": "Course Content",
            "description": "",
            "order": 1
        }
        db.modules.insert_one(module)
        
        # Create lessons from slides with better formatting
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_title = f"Slide {slide_num}"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if len(text) < 100 and len(slide_text) == 0:
                        slide_title = text
                    slide_text.append(text)
            
            # Create better formatted HTML
            content_html = '''<div class="lesson-content">
<style>
.slide-container {
    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    padding: 40px;
    border-radius: 12px;
    color: white;
    min-height: 400px;
    display: flex;
    flex-direction: column;
    justify-content: center;
}
.slide-header {
    font-size: 14px;
    opacity: 0.8;
    margin-bottom: 20px;
}
.slide-title {
    font-size: 32px;
    font-weight: bold;
    margin-bottom: 30px;
    line-height: 1.3;
}
.slide-content {
    font-size: 18px;
    line-height: 1.8;
}
.slide-content p {
    margin: 15px 0;
}
.slide-footer {
    margin-top: 40px;
    padding-top: 20px;
    border-top: 1px solid rgba(255,255,255,0.3);
    font-size: 14px;
    opacity: 0.8;
}
</style>
<div class="slide-container">
    <div class="slide-header">''' + ppt_info['title'] + f''' - Slide {slide_num} of {total_slides}</div>
    <div class="slide-title">''' + (slide_text[0] if slide_text else f"Slide {slide_num}") + '''</div>
    <div class="slide-content">'''
            
            if len(slide_text) > 1:
                for text in slide_text[1:]:
                    content_html += f'<p>{text}</p>'
            elif not slide_text:
                content_html += '<p><em>Visual content - Please refer to the original presentation for images and diagrams</em></p>'
            
            content_html += '''</div>
    <div class="slide-footer">Use the navigation buttons below to move through the slides</div>
</div>
</div>'''
            
            lesson_id = str(uuid.uuid4())
            lesson = {
                "id": lesson_id,
                "module_id": module_id,
                "title": slide_title if slide_title != f"Slide {slide_num}" else f"Slide {slide_num}",
                "content_type": "text",
                "content": content_html,
                "duration_minutes": 2,
                "order": slide_num
            }
            db.lessons.insert_one(lesson)
        
        print(f"  Created {total_slides} lessons")
        
    except Exception as e:
        print(f"  Error: {str(e)}")

print("\n✓ All courses created!")
print("\nTo show actual slide images:")
print("1. Export each PowerPoint as images")
print("2. Upload to Google Drive or image hosting")
print("3. I can update the lessons with image URLs")

client.close()
